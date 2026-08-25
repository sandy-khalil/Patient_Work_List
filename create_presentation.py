from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Colors ──────────────────────────────────────────────
PRIMARY   = RGBColor(0x1B, 0x4F, 0x72)   # dark blue
ACCENT    = RGBColor(0x2E, 0x86, 0xC1)   # bright blue
LIGHT_BG  = RGBColor(0xEB, 0xF5, 0xFB)   # light blue bg
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
BLACK     = RGBColor(0x2C, 0x3E, 0x50)
GRAY      = RGBColor(0x7F, 0x8C, 0x8D)
GREEN     = RGBColor(0x27, 0xAE, 0x60)
ORANGE    = RGBColor(0xE6, 0x7E, 0x22)
RED       = RGBColor(0xC0, 0x39, 0x2B)
DARK_BG   = RGBColor(0x1A, 0x25, 0x2F)


def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_bg(slide, x, y, w, h, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if alpha is not None:
        from lxml import etree
        nsmap = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
        solidFill = shape.fill._fill
        srgb = solidFill.find('.//a:srgbClr', nsmap)
        if srgb is None:
            srgb = solidFill.find('.//a:srgbClr', nsmap)
        if srgb is not None:
            alpha_el = etree.SubElement(srgb, '{http://schemas.openxmlformats.org/drawingml/2006/main}alpha')
            alpha_el.set('val', str(int(alpha * 1000)))
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=BLACK, bold=False, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=BLACK, spacing=Pt(8), bullet_char='●'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"{bullet_char}  {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Calibri'
        p.space_after = spacing
        p.space_before = Pt(2)
    return txBox


def add_code_block(slide, left, top, width, height, code_text, font_size=13):
    shape = add_shape_bg(slide, left, top, width, height, RGBColor(0x1E, 0x1E, 0x2E))
    txBox = slide.shapes.add_textbox(left + Inches(0.25), top + Inches(0.2),
                                      width - Inches(0.5), height - Inches(0.4))
    tf = txBox.text_frame
    tf.word_wrap = True
    lines = code_text.strip().split('\n')
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = RGBColor(0xA6, 0xE2, 0x2E)
        p.font.name = 'Consolas'
        p.space_after = Pt(1)
    return shape


def slide_number_footer(slide, num, total=10):
    add_text_box(slide, Inches(12.2), Inches(7.05), Inches(1), Inches(0.4),
                 f"{num} / {total}", font_size=10, color=GRAY, alignment=PP_ALIGN.RIGHT)


def add_icon_box(slide, left, top, icon_text, label, color=ACCENT):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(1.8), Inches(1.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = color
    shape.line.width = Pt(2)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    p.text = icon_text
    p.font.size = Pt(28)
    p.font.color.rgb = color
    p.font.bold = True
    p.font.name = 'Calibri'
    p2 = tf.add_paragraph()
    p2.text = label
    p2.font.size = Pt(12)
    p2.font.color.rgb = BLACK
    p2.font.name = 'Calibri'
    p2.alignment = PP_ALIGN.CENTER
    return shape


# ═══════════════════════════════════════════════════════
# SLIDE 1 — Title Slide
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, DARK_BG)

# accent bar
add_shape_bg(slide, Inches(0), Inches(2.8), Inches(13.333), Inches(0.06), ACCENT)

add_text_box(slide, Inches(1.5), Inches(1.2), Inches(10), Inches(1.2),
             "PATIENT WORKLIST", font_size=48, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER, font_name='Calibri')

add_text_box(slide, Inches(1.5), Inches(2.0), Inches(10), Inches(0.7),
             "Hospital Imaging Study Management System", font_size=22, color=ACCENT,
             alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(3.5), Inches(10), Inches(0.7),
             "Full-Stack Web Application  |  ASP.NET Core  |  Entity Framework  |  SQLite",
             font_size=16, color=GRAY, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(5.5), Inches(10), Inches(0.6),
             "Sandy Khalil  |  Millensys Final Project  |  2026",
             font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)

slide_number_footer(slide, 1)


# ═══════════════════════════════════════════════════════
# SLIDE 2 — Agenda / Table of Contents
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

# header bar
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.1), PRIMARY)
add_text_box(slide, Inches(0.8), Inches(0.2), Inches(8), Inches(0.7),
             "AGENDA", font_size=36, color=WHITE, bold=True)

items_left = [
    "1.  Project Overview",
    "2.  Technology Stack",
    "3.  Data Model & Architecture",
    "4.  API Endpoints",
    "5.  Frontend Overview",
]
items_right = [
    "6.  Database Design",
    "7.  Repository Pattern",
    "8.  Security & Validation",
    "9.  Demo Walkthrough",
    "10. Summary & Future Work",
]

add_bullet_list(slide, Inches(1), Inches(1.8), Inches(5), Inches(4.5),
                items_left, font_size=20, color=BLACK, bullet_char='▸')
add_bullet_list(slide, Inches(7), Inches(1.8), Inches(5), Inches(4.5),
                items_right, font_size=20, color=BLACK, bullet_char='▸')

slide_number_footer(slide, 2)


# ═══════════════════════════════════════════════════════
# SLIDE 3 — Project Overview
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.1), PRIMARY)
add_text_box(slide, Inches(0.8), Inches(0.2), Inches(8), Inches(0.7),
             "PROJECT OVERVIEW", font_size=36, color=WHITE, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.7),
             "A hospital-style web application for managing patients, doctors, and imaging studies.",
             font_size=20, color=BLACK)

# feature boxes
features = [
    ("👥", "Patients", "Demographics, MRN, status tracking"),
    ("🩺", "Doctors", "Specialty management, linked records"),
    ("📋", "Studies", "CT, MRI, X-Ray scheduling & status"),
    ("📊", "Worklist", "Real-time overview of all studies"),
]

for i, (icon, title, desc) in enumerate(features):
    left = Inches(0.8 + i * 3.1)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(2.6),
                                    Inches(2.8), Inches(2.0))
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_BG
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    p.text = icon
    p.font.size = Pt(36)
    p2 = tf.add_paragraph()
    p2.text = title
    p2.font.size = Pt(18)
    p2.font.bold = True
    p2.font.color.rgb = PRIMARY
    p2.alignment = PP_ALIGN.CENTER
    p3 = tf.add_paragraph()
    p3.text = desc
    p3.font.size = Pt(12)
    p3.font.color.rgb = GRAY
    p3.alignment = PP_ALIGN.CENTER

add_text_box(slide, Inches(0.8), Inches(5.2), Inches(11.5), Inches(1.5),
             "The application simulates a radiology department worklist — the screen used to track\n"
             "which imaging studies are scheduled, in progress, or completed for each patient.",
             font_size=16, color=GRAY)

slide_number_footer(slide, 3)


# ═══════════════════════════════════════════════════════
# SLIDE 4 — Technology Stack
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.1), PRIMARY)
add_text_box(slide, Inches(0.8), Inches(0.2), Inches(8), Inches(0.7),
             "TECHNOLOGY STACK", font_size=36, color=WHITE, bold=True)

# Table-style layout
cols = [
    ("Backend", ["ASP.NET Core 7", "C# (.NET 7)", "Minimal API hosting", "Swagger / Swashbuckle"], ACCENT),
    ("Data Layer", ["Entity Framework Core 7", "SQLite Database", "Repository Pattern", "LINQ Queries"], GREEN),
    ("Frontend", ["HTML5 / CSS3", "Bootstrap 5.3", "jQuery 3.7", "DataTables 1.13"], ORANGE),
    ("Architecture", ["RESTful API", "DTO Pattern", "Dependency Injection", "Layered Architecture"], RED),
]

for i, (title, items, color) in enumerate(cols):
    left = Inches(0.6 + i * 3.2)
    # column header
    header = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.5),
                                     Inches(2.9), Inches(0.6))
    header.fill.solid()
    header.fill.fore_color.rgb = color
    header.line.fill.background()
    tf = header.text_frame
    tf.paragraphs[0].text = title
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    for j, item in enumerate(items):
        y = Inches(2.3 + j * 0.7)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, y,
                                      Inches(2.9), Inches(0.55))
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_BG
        box.line.fill.background()
        tf = box.text_frame
        tf.paragraphs[0].text = item
        tf.paragraphs[0].font.size = Pt(13)
        tf.paragraphs[0].font.color.rgb = BLACK
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

slide_number_footer(slide, 4)


# ═══════════════════════════════════════════════════════
# SLIDE 5 — Data Model
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.1), PRIMARY)
add_text_box(slide, Inches(0.8), Inches(0.2), Inches(8), Inches(0.7),
             "DATA MODEL", font_size=36, color=WHITE, bold=True)

# Entity boxes
entities = [
    ("Person", "PersonId (PK)\nFirstName\nLastName\nDateOfBirth\nGender\nPhone?\nEmail?", Inches(0.5), ACCENT),
    ("Patient", "PatientId (PK)\nPersonId (FK)\nMRN (unique)\nStatus", Inches(3.8), GREEN),
    ("Doctor", "DoctorId (PK)\nPersonId (FK)\nSpecialty", Inches(7.1), ORANGE),
    ("Study", "StudyId (PK)\nPatientId (FK)\nDoctorId (FK)\nModality\nStudyDate\nStatus", Inches(10.2), RED),
]

for name, fields, left, color in entities:
    # entity header
    header = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.5),
                                     Inches(2.7), Inches(0.5))
    header.fill.solid()
    header.fill.fore_color.rgb = color
    header.line.fill.background()
    tf = header.text_frame
    tf.paragraphs[0].text = name
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # entity fields
    body = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(2.05),
                                   Inches(2.7), Inches(2.4))
    body.fill.solid()
    body.fill.fore_color.rgb = LIGHT_BG
    body.line.color.rgb = color
    body.line.width = Pt(1.5)
    tf = body.text_frame
    tf.word_wrap = True
    lines = fields.split('\n')
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(12)
        p.font.color.rgb = BLACK
        p.font.name = 'Consolas'
        p.space_after = Pt(3)

# Relationship annotations
add_text_box(slide, Inches(0.5), Inches(4.8), Inches(12), Inches(0.5),
             "Relationships:  Person 1:1 Patient  |  Person 1:1 Doctor  |  Patient 1:N Studies  |  Doctor 1:N Studies",
             font_size=15, color=PRIMARY, bold=True, alignment=PP_ALIGN.CENTER)

# Design note
note_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.5),
                                     Inches(12.3), Inches(1.3))
note_shape.fill.solid()
note_shape.fill.fore_color.rgb = RGBColor(0xFD, 0xF2, 0xE9)
note_shape.line.color.rgb = ORANGE
note_shape.line.width = Pt(1)
tf = note_shape.text_frame
tf.word_wrap = True
tf.paragraphs[0].text = "💡 Design Decision: Why a shared Person table?"
tf.paragraphs[0].font.size = Pt(14)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = ORANGE
p2 = tf.add_paragraph()
p2.text = ("Patients and Doctors share the same demographic fields (name, DOB, gender, contact). "
           "A shared Person table eliminates duplication and allows a person to be both a patient and a doctor.")
p2.font.size = Pt(12)
p2.font.color.rgb = BLACK

slide_number_footer(slide, 5)


# ═══════════════════════════════════════════════════════
# SLIDE 6 — Architecture
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.1), PRIMARY)
add_text_box(slide, Inches(0.8), Inches(0.2), Inches(8), Inches(0.7),
             "APPLICATION ARCHITECTURE", font_size=36, color=WHITE, bold=True)

# Layer boxes (stacked vertically)
layers = [
    ("Frontend (wwwroot)", "index.html  •  app.js  •  style.css", "Bootstrap 5, jQuery, DataTables", LIGHT_BG, ACCENT),
    ("Controllers", "PatientsController  •  DoctorsController  •  StudiesController", "HTTP requests, validation, DTO mapping", RGBColor(0xE8, 0xF8, 0xF5), GREEN),
    ("Repositories", "GenericRepository<T>  •  PatientRepository  •  DoctorRepository  •  StudyRepository", "Data access, eager loading, CRUD", RGBColor(0xFD, 0xF2, 0xE9), ORANGE),
    ("EF Core / SQLite", "ApplicationDbContext  •  Entities  •  DbSeeder", "ORM, change tracking, migrations", RGBColor(0xFD, 0xED, 0xED), RED),
]

for i, (title, details, tech, bg_color, accent) in enumerate(layers):
    y = Inches(1.4 + i * 1.45)
    # layer box
    layer = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), y,
                                    Inches(10.3), Inches(1.25))
    layer.fill.solid()
    layer.fill.fore_color.rgb = bg_color
    layer.line.color.rgb = accent
    layer.line.width = Pt(2)

    tf = layer.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = title
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = accent
    p2 = tf.add_paragraph()
    p2.text = details
    p2.font.size = Pt(12)
    p2.font.color.rgb = BLACK
    p2.font.name = 'Consolas'
    p3 = tf.add_paragraph()
    p3.text = tech
    p3.font.size = Pt(11)
    p3.font.color.rgb = GRAY

    # arrow between layers
    if i < len(layers) - 1:
        add_text_box(slide, Inches(6.2), y + Inches(1.25), Inches(1), Inches(0.3),
                     "▼", font_size=18, color=GRAY, alignment=PP_ALIGN.CENTER)

slide_number_footer(slide, 6)


# ═══════════════════════════════════════════════════════
# SLIDE 7 — API Endpoints
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.1), PRIMARY)
add_text_box(slide, Inches(0.8), Inches(0.2), Inches(8), Inches(0.7),
             "REST API ENDPOINTS", font_size=36, color=WHITE, bold=True)

api_data = [
    ("GET", "/api/patients", "List all patients with details"),
    ("POST", "/api/patients", "Create a new patient"),
    ("GET", "/api/patients/{id}", "Get patient by ID"),
    ("PUT", "/api/patients/{id}", "Update patient"),
    ("DELETE", "/api/patients/{id}", "Delete patient + studies"),
    ("GET", "/api/doctors", "List all doctors"),
    ("POST", "/api/doctors", "Create a new doctor"),
    ("DELETE", "/api/doctors/{id}", "Delete (only if no studies)"),
    ("GET", "/api/studies", "List studies (filterable)"),
    ("POST", "/api/studies", "Create study (validates FK)"),
]

method_colors = {
    'GET': GREEN,
    'POST': ACCENT,
    'PUT': ORANGE,
    'DELETE': RED,
}

# Table header
header_y = Inches(1.35)
add_shape_bg(slide, Inches(0.8), header_y, Inches(11.7), Inches(0.45), PRIMARY)
for col_x, col_w, col_text in [
    (Inches(0.8), Inches(1.2), "Method"),
    (Inches(2.0), Inches(4.0), "Endpoint"),
    (Inches(6.0), Inches(6.5), "Description"),
]:
    add_text_box(slide, col_x + Inches(0.15), header_y + Inches(0.05), col_w, Inches(0.35),
                 col_text, font_size=13, color=WHITE, bold=True)

for i, (method, endpoint, desc) in enumerate(api_data):
    y = Inches(1.85 + i * 0.48)
    bg = LIGHT_BG if i % 2 == 0 else WHITE
    add_shape_bg(slide, Inches(0.8), y, Inches(11.7), Inches(0.45), bg)

    # method badge
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.95), y + Inches(0.05),
                                    Inches(1.0), Inches(0.35))
    badge.fill.solid()
    badge.fill.fore_color.rgb = method_colors.get(method, GRAY)
    badge.line.fill.background()
    tf = badge.text_frame
    tf.paragraphs[0].text = method
    tf.paragraphs[0].font.size = Pt(10)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_text_box(slide, Inches(2.1), y + Inches(0.05), Inches(3.8), Inches(0.35),
                 endpoint, font_size=12, color=BLACK, font_name='Consolas')
    add_text_box(slide, Inches(6.1), y + Inches(0.05), Inches(6.3), Inches(0.35),
                 desc, font_size=12, color=GRAY)

slide_number_footer(slide, 7)


# ═══════════════════════════════════════════════════════
# SLIDE 8 — Repository Pattern
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.1), PRIMARY)
add_text_box(slide, Inches(0.8), Inches(0.2), Inches(8), Inches(0.7),
             "REPOSITORY PATTERN", font_size=36, color=WHITE, bold=True)

# Generic interface code
add_text_box(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(0.4),
             "Generic Repository Interface", font_size=16, color=PRIMARY, bold=True)

code_generic = """public interface IRepository<T> where T : class
{
    Task<IEnumerable<T>> GetAllAsync();
    Task<T?> GetByIdAsync(int id);
    Task<T> AddAsync(T entity);
    Task UpdateAsync(T entity);
    Task DeleteAsync(T entity);
    Task<bool> ExistsAsync(int id);
}"""

add_code_block(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(3.2), code_generic, font_size=13)

# Specific repository example
add_text_box(slide, Inches(7), Inches(1.3), Inches(5.5), Inches(0.4),
             "Specific Repository (Patient)", font_size=16, color=GREEN, bold=True)

code_specific = """public class PatientRepository
    : GenericRepository<Patient>,
      IPatientRepository
{
    public async Task<IEnumerable<Patient>>
        GetAllWithDetailsAsync()
    {
        return await _dbSet
            .Include(p => p.Person)
            .Include(p => p.Studies)
            .AsNoTracking()
            .ToListAsync();
    }
}"""

add_code_block(slide, Inches(7), Inches(1.8), Inches(5.5), Inches(3.2), code_specific, font_size=12)

# Benefits box
benefits = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.3),
                                   Inches(11.7), Inches(1.6))
benefits.fill.solid()
benefits.fill.fore_color.rgb = RGBColor(0xE8, 0xF8, 0xF5)
benefits.line.color.rgb = GREEN
benefits.line.width = Pt(1.5)
tf = benefits.text_frame
tf.word_wrap = True
tf.paragraphs[0].text = "Key Benefits"
tf.paragraphs[0].font.size = Pt(16)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = GREEN

benefits_list = [
    "Separation of concerns — controllers never touch DbContext directly",
    "Testability — repositories can be mocked for unit tests",
    "Reusability — generic CRUD shared across all entities",
    "Specific repositories add eager loading (.Include) and custom queries",
]
for item in benefits_list:
    p = tf.add_paragraph()
    p.text = f"▸  {item}"
    p.font.size = Pt(12)
    p.font.color.rgb = BLACK
    p.space_after = Pt(2)

slide_number_footer(slide, 8)


# ═══════════════════════════════════════════════════════
# SLIDE 9 — Frontend Overview
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.1), PRIMARY)
add_text_box(slide, Inches(0.8), Inches(0.2), Inches(8), Inches(0.7),
             "FRONTEND OVERVIEW", font_size=36, color=WHITE, bold=True)

# Three tab mockups
tabs = [
    ("Studies Tab", "View all imaging studies\nFilter by patient/doctor\nModality, date, status\nCreate / Edit / Delete", ACCENT),
    ("Patients Tab", "Patient demographics\nMRN & status\nAge calculator\nStudy count per patient", GREEN),
    ("Doctors Tab", "Doctor demographics\nSpecialty management\nLinked study count\nDelete protection", ORANGE),
]

for i, (title, features, color) in enumerate(tabs):
    left = Inches(0.8 + i * 4.2)

    # tab header
    tab_h = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.5),
                                    Inches(3.8), Inches(0.5))
    tab_h.fill.solid()
    tab_h.fill.fore_color.rgb = color
    tab_h.line.fill.background()
    tf = tab_h.text_frame
    tf.paragraphs[0].text = title
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # tab body
    tab_b = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(2.05),
                                    Inches(3.8), Inches(2.2))
    tab_b.fill.solid()
    tab_b.fill.fore_color.rgb = LIGHT_BG
    tab_b.line.color.rgb = color
    tab_b.line.width = Pt(1.5)
    tf = tab_b.text_frame
    tf.word_wrap = True
    for j, line in enumerate(features.split('\n')):
        if j == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"▸  {line}"
        p.font.size = Pt(13)
        p.font.color.rgb = BLACK
        p.space_after = Pt(6)

# Tech stack footer
add_shape_bg(slide, Inches(0.8), Inches(4.8), Inches(11.7), Inches(2.0), LIGHT_BG)
add_text_box(slide, Inches(1.2), Inches(4.9), Inches(11), Inches(0.4),
             "Frontend Technologies", font_size=16, color=PRIMARY, bold=True)

techs = [
    ("Bootstrap 5.3", "Responsive grid, modals, toasts, badges"),
    ("jQuery 3.7", "DOM manipulation, event handling"),
    ("DataTables 1.13", "Sortable/searchable/paginated tables"),
    ("Fetch API", "RESTful calls to backend, error handling"),
]
for i, (name, desc) in enumerate(techs):
    col = Inches(1.2 + (i % 2) * 6)
    row = Inches(5.4 + (i // 2) * 0.55)
    add_text_box(slide, col, row, Inches(2), Inches(0.3),
                 name, font_size=13, color=ACCENT, bold=True)
    add_text_box(slide, col + Inches(2.1), row, Inches(3.5), Inches(0.3),
                 desc, font_size=12, color=GRAY)

slide_number_footer(slide, 9)


# ═══════════════════════════════════════════════════════
# SLIDE 10 — Summary & Future Work
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)

add_shape_bg(slide, Inches(0), Inches(0.8), Inches(13.333), Inches(0.06), ACCENT)

add_text_box(slide, Inches(1), Inches(0.15), Inches(11), Inches(0.7),
             "SUMMARY & FUTURE WORK", font_size=36, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)

# Two columns
add_text_box(slide, Inches(1), Inches(1.3), Inches(5.5), Inches(0.5),
             "What Was Built", font_size=20, color=ACCENT, bold=True)

built_items = [
    "Full-stack hospital worklist application",
    "RESTful API with 15 endpoints",
    "SQLite database with seeded demo data",
    "Repository pattern for clean architecture",
    "Responsive Bootstrap frontend",
    "Swagger UI for API documentation",
]
add_bullet_list(slide, Inches(1), Inches(1.9), Inches(5.5), Inches(3.5),
                built_items, font_size=14, color=RGBColor(0xEC, 0xF0, 0xF1), bullet_char='✓')

add_text_box(slide, Inches(7), Inches(1.3), Inches(5.5), Inches(0.5),
             "Future Enhancements", font_size=20, color=ORANGE, bold=True)

future_items = [
    "JWT authentication & role-based access",
    "PostgreSQL / SQL Server migration",
    "EF Core Migrations (instead of EnsureCreated)",
    "Unit & integration tests",
    "Docker containerization",
    "Image upload for study reports",
]
add_bullet_list(slide, Inches(7), Inches(1.9), Inches(5.5), Inches(3.5),
                future_items, font_size=14, color=RGBColor(0xEC, 0xF0, 0xF1), bullet_char='→')

# footer
add_shape_bg(slide, Inches(0), Inches(5.8), Inches(13.333), Inches(0.06), ACCENT)
add_text_box(slide, Inches(1), Inches(6.1), Inches(11.3), Inches(0.5),
             "Thank You  |  Sandy Khalil  |  Millensys Final Project",
             font_size=16, color=GRAY, alignment=PP_ALIGN.CENTER)

slide_number_footer(slide, 10)


# ── Save ───────────────────────────────────────────────
output_path = "/media/sandy/Work2/millensys/Final Project/Patient_Worklist_Presentation.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
